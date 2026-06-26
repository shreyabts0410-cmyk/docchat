"""Document extraction & editing for docx, pdf, pptx, xlsx."""
import io
import re
from typing import List, Tuple

from docx import Document as DocxDocument
from docx.oxml.ns import qn
from pypdf import PdfReader
from pptx import Presentation
from openpyxl import load_workbook


# ── extraction helpers ──────────────────────────────────────────────

def _extract_docx(b: bytes) -> str:
    doc = DocxDocument(io.BytesIO(b))
    parts = []
    # Extract paragraphs
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            row_vals = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_vals:
                parts.append(" | ".join(row_vals))
    return "\n".join(parts)


def _extract_pdf(b: bytes) -> str:
    reader = PdfReader(io.BytesIO(b))
    parts: List[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _extract_pptx(b: bytes) -> str:
    prs = Presentation(io.BytesIO(b))
    parts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _extract_xlsx(b: bytes) -> str:
    wb = load_workbook(io.BytesIO(b), data_only=True)
    parts: List[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            vals = [str(c.value) for c in row if c.value is not None]
            if vals:
                parts.append("\t".join(vals))
    return "\n".join(parts)


_EXTRACTORS = {
    "docx": _extract_docx,
    "pdf": _extract_pdf,
    "pptx": _extract_pptx,
    "xlsx": _extract_xlsx,
}


def extract_text(file_bytes: bytes, ext: str) -> str:
    """Return plain-text content of a document."""
    ext = ext.lower().lstrip(".")
    fn = _EXTRACTORS.get(ext)
    if fn is None:
        raise ValueError(f"Unsupported extension: {ext}")
    return fn(file_bytes)


# ── DOCX editing helpers ────────────────────────────────────────────

def _para_full_text(para) -> str:
    """Get the full text of a paragraph across all runs."""
    return "".join(run.text for run in para.runs)


def _copy_run_format(src_run, dst_run):
    """Copy font formatting from src_run to dst_run."""
    try:
        dst_run.bold = src_run.bold
        dst_run.italic = src_run.italic
        dst_run.underline = src_run.underline
        if src_run.font.name:
            dst_run.font.name = src_run.font.name
        if src_run.font.size:
            dst_run.font.size = src_run.font.size
        if src_run.font.color and src_run.font.color.type:
            dst_run.font.color.rgb = src_run.font.color.rgb
    except Exception:
        pass


def _replace_in_para(para, old: str, new: str) -> bool:
    """
    Replace old text with new text inside a paragraph.
    Preserves per-run formatting (font, bold, italic, size, color).

    Strategy:
    1. Check if old text exists across the full paragraph text
    2. Find which run(s) contain the old text
    3. If it fits within a single run — replace just that run's text, all formatting preserved
    4. If it spans multiple runs — do a careful rebuild keeping each run's format
    """
    full_text = _para_full_text(para)
    if old not in full_text:
        return False

    # Case 1: old text fits entirely within a single run — cleanest case
    for run in para.runs:
        if old in run.text:
            run.text = run.text.replace(old, new, 1)
            return True

    # Case 2: old text spans multiple runs — need to rebuild carefully
    # Find where old starts and ends in the full text
    start_idx = full_text.index(old)
    end_idx = start_idx + len(old)

    # Map character positions to runs
    run_ranges = []
    pos = 0
    for run in para.runs:
        run_ranges.append((pos, pos + len(run.text), run))
        pos += len(run.text)

    # Find which runs are involved in the match
    affected = []
    for r_start, r_end, run in run_ranges:
        if r_end > start_idx and r_start < end_idx:
            affected.append((r_start, r_end, run))

    if not affected:
        return False

    # Keep the formatting of the first affected run for the new text
    first_run = affected[0][2]

    # Rebuild: for each affected run, figure out what part to keep
    for r_start, r_end, run in affected:
        # Part before the match in this run
        before = ""
        if r_start < start_idx:
            before = run.text[:start_idx - r_start]

        # Part after the match in this run
        after = ""
        if r_end > end_idx:
            after = run.text[end_idx - r_start:]

        if run is first_run.text and run is affected[0][2]:
            # First affected run: keep before + new text
            run.text = before + new
            _copy_run_format(first_run, run)
        elif run is affected[-1][2] and run is not affected[0][2]:
            # Last affected run: keep only after
            run.text = after
        else:
            # Middle runs: clear completely
            run.text = ""

    # Simpler fallback: if the above left things messy, just do the safe version
    # Re-check if replacement actually worked
    if old not in _para_full_text(para) and new in _para_full_text(para):
        return True

    # Final fallback: collapse into first run with its formatting preserved
    new_full = full_text.replace(old, new, 1)
    if para.runs:
        first = para.runs[0]
        saved_bold = first.bold
        saved_italic = first.italic
        saved_underline = first.underline
        saved_font_name = first.font.name
        saved_font_size = first.font.size
        first.text = new_full
        first.bold = saved_bold
        first.italic = saved_italic
        first.underline = saved_underline
        if saved_font_name:
            first.font.name = saved_font_name
        if saved_font_size:
            first.font.size = saved_font_size
        for run in para.runs[1:]:
            run.text = ""
    return True


def _replace_in_para_append(para, anchor: str, append_text: str) -> bool:
    """Append text after the anchor text in a paragraph, preserving formatting."""
    full_text = _para_full_text(para)
    if anchor not in full_text:
        return False
    return _replace_in_para(para, anchor, anchor + append_text)


def _edit_docx(b: bytes, old: str, new: str) -> Tuple[bytes, bool]:
    doc = DocxDocument(io.BytesIO(b))
    changed = False

    # 1. Search in body paragraphs
    for para in doc.paragraphs:
        if _replace_in_para(para, old, new):
            changed = True
            break

    # 2. Search in tables (covers forms, school forms, data tables)
    if not changed:
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        if _replace_in_para(para, old, new):
                            changed = True
                            break
                    if changed:
                        break
                if changed:
                    break
            if changed:
                break

    # 3. Search in text boxes and other shapes
    if not changed:
        for shape in doc.element.body.iter():
            tag = shape.tag.split("}")[-1] if "}" in shape.tag else shape.tag
            if tag == "txbx":
                for para_elem in shape.iter():
                    para_tag = para_elem.tag.split("}")[-1] if "}" in para_elem.tag else para_elem.tag
                    if para_tag == "p":
                        from docx.text.paragraph import Paragraph
                        para_obj = Paragraph(para_elem, doc)
                        if _replace_in_para(para_obj, old, new):
                            changed = True
                            break
                if changed:
                    break

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue(), changed


def _edit_pptx(b: bytes, old: str, new: str) -> Tuple[bytes, bool]:
    prs = Presentation(io.BytesIO(b))
    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = "".join(run.text for run in para.runs)
                    if old in full_text:
                        new_full = full_text.replace(old, new, 1)
                        if para.runs:
                            para.runs[0].text = new_full
                            for run in para.runs[1:]:
                                run.text = ""
                        changed = True
                        break
            if changed:
                break
        if changed:
            break
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), changed


def _edit_xlsx(b: bytes, old: str, new: str) -> Tuple[bytes, bool]:
    wb = load_workbook(io.BytesIO(b))
    changed = False
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None and isinstance(cell.value, str) and old in cell.value:
                    cell.value = cell.value.replace(old, new, 1)
                    changed = True
                    break
            if changed:
                break
        if changed:
            break
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue(), changed


def _edit_pdf(b: bytes, old: str, new: str) -> Tuple[bytes, bool]:
    try:
        import fitz  # pymupdf
    except ImportError:
        return b, False

    doc = fitz.open(stream=b, filetype="pdf")
    changed = False

    for page in doc:
        # Search for the text on this page
        instances = page.search_for(old)
        if not instances:
            continue

        # Take the first match only
        rect = instances[0]

        # Get the font properties of the existing text at that location
        blocks = page.get_text("dict")["blocks"]
        font_size = 11  # default
        font_name = "helv"  # default
        font_color = (0, 0, 0)  # default black

        for block in blocks:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    span_rect = fitz.Rect(span["bbox"])
                    if span_rect.intersects(rect):
                        font_size = span["size"]
                        font_color_int = span.get("color", 0)
                        # Convert int color to RGB tuple (0-1 range)
                        r = ((font_color_int >> 16) & 0xFF) / 255
                        g = ((font_color_int >> 8) & 0xFF) / 255
                        b_val = (font_color_int & 0xFF) / 255
                        font_color = (r, g, b_val)
                        break

        # Redact (white out) the old text area
        page.add_redact_annot(rect, fill=(1, 1, 1))
        page.apply_redactions()

        # Insert new text in the same position with same font properties
        page.insert_text(
            rect.tl,  # top-left point
            new,
            fontsize=font_size,
            color=font_color,
        )
        changed = True
        break  # only first occurrence

    if changed:
        out = io.BytesIO()
        doc.save(out)
        doc.close()
        return out.getvalue(), True

    doc.close()
    return b, False


_EDITORS = {
    "docx": _edit_docx,
    "pdf": _edit_pdf,
    "pptx": _edit_pptx,
    "xlsx": _edit_xlsx,
}


def apply_edit(file_bytes: bytes, ext: str, old: str, new: str) -> Tuple[bytes, bool]:
    """Replace first occurrence of *old* with *new* in the document. Returns (new_bytes, changed)."""
    ext = ext.lower().lstrip(".")
    fn = _EDITORS.get(ext)
    if fn is None:
        raise ValueError(f"Unsupported extension: {ext}")
    return fn(file_bytes, old, new)
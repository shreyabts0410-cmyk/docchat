from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import os
import shutil
import json
import re
import uuid
from typing import List, Dict, Optional

from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader, CSVLoader
from langchain.docstore.document import Document
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.messages import SystemMessage, HumanMessage, AIMessage
from dotenv import load_dotenv
import pandas as pd
import google.generativeai as genai
from PIL import Image

load_dotenv()

genai.configure(api_key=os.environ.get("GOOGLE_API_KEY"))
image_model = genai.GenerativeModel('gemini-3.1-flash-lite')

app = FastAPI(title="AI Document Analysis Platform")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Mount static files to serve documents to the frontend viewer
app.mount("/files", StaticFiles(directory=UPLOAD_DIR), name="files")

# In-memory storage for uploaded documents' full text (kept for insights generation)
uploaded_documents: List[Dict[str, str]] = []

# In-memory storage for structured content maps (keyed by file_id)
content_maps: Dict[str, Dict] = {}

# ──────────────────────────────────────────────────────────────
# TONE ADAPTATION ENGINE
# ──────────────────────────────────────────────────────────────

DOCUMENT_STYLE_PROFILES = {
    "legal": {
        "label": "Legal / Regulatory",
        "tone": "formal, precise, and authoritative",
        "directives": (
            "Use precise legal terminology. Cite specific clauses, sections, or provisions when referencing document content. "
            "Maintain a measured, objective tone. Avoid colloquialisms. Structure responses with numbered points when listing obligations or risks. "
            "When editing, preserve legal formality and ensure no ambiguity is introduced."
        ),
    },
    "technical": {
        "label": "Technical / Engineering",
        "tone": "clear, methodical, and technically rigorous",
        "directives": (
            "Use domain-appropriate technical vocabulary. Be concise and structured — prefer bullet points, tables, and step-by-step breakdowns. "
            "Reference specifications, metrics, and data points when available. When editing, prioritize precision and eliminate vague qualifiers."
        ),
    },
    "medical": {
        "label": "Medical / Clinical",
        "tone": "clinical, evidence-based, and empathetic when appropriate",
        "directives": (
            "Use standard medical nomenclature (ICD, SNOMED, drug names). Distinguish clearly between findings, diagnoses, and recommendations. "
            "Flag safety-critical information prominently. When editing, ensure clinical accuracy and appropriate hedging for uncertain findings."
        ),
    },
    "academic": {
        "label": "Academic / Research",
        "tone": "scholarly, analytical, and balanced",
        "directives": (
            "Use academic conventions: hedge appropriately, cite evidence, distinguish between correlation and causation. "
            "Structure analysis around hypotheses, methodology, findings, and implications. When editing, elevate to publication-ready prose."
        ),
    },
    "financial": {
        "label": "Financial / Business",
        "tone": "executive-grade, data-driven, and action-oriented",
        "directives": (
            "Lead with bottom-line impact. Reference specific figures, KPIs, and financial metrics. "
            "Use concise executive language — no filler. Structure around key takeaways, risks, and recommended actions. "
            "When editing, sharpen for C-suite readability."
        ),
    },
    "marketing": {
        "label": "Marketing / Creative",
        "tone": "engaging, persuasive, and brand-aware",
        "directives": (
            "Match the brand voice present in the document. Use compelling, audience-centric language. "
            "When editing, strengthen hooks, CTAs, and emotional resonance while maintaining authenticity."
        ),
    },
    "casual": {
        "label": "General / Informal",
        "tone": "friendly, clear, and conversational",
        "directives": (
            "Use approachable language without being unprofessional. Explain concepts simply. "
            "When editing, prioritize readability and natural flow over formality."
        ),
    },
}

INTENT_TONE_MODIFIERS = {
    "analysis": "Provide deep, structured analysis. Use evidence from the document to support every point. Be thorough.",
    "rewrite": "Focus on producing a polished, improved version of the text. Show the rewritten version clearly so it can be copied directly.",
    "summarize": "Be concise and distill to essentials. Lead with the most critical information. Use bullet points for scanability.",
    "compare": "Structure the comparison clearly — use tables or side-by-side points. Highlight key differences and similarities.",
    "extract": "Pull out the specific data points, facts, or clauses requested. Present them clearly and cite their location in the document.",
    "casual": "Be conversational and direct. Answer naturally without excessive formality.",
}


def classify_document_style(text: str) -> str:
    """Classify the style/domain of a document based on content heuristics."""
    text_lower = text[:10000].lower()

    scores: Dict[str, int] = {style: 0 for style in DOCUMENT_STYLE_PROFILES}

    # Legal signals
    legal_terms = ["whereas", "hereinafter", "indemnif", "arbitrat", "jurisdiction", "clause", "statute",
                   "liability", "breach", "contract", "plaintiff", "defendant", "pursuant", "herein",
                   "enforceable", "stipulat", "affidavit", "tort", "negligence", "compliance"]
    scores["legal"] += sum(3 for t in legal_terms if t in text_lower)

    # Technical signals
    tech_terms = ["api", "algorithm", "architecture", "deploy", "infrastructure", "latency",
                  "throughput", "configuration", "protocol", "endpoint", "microservice", "docker",
                  "kubernetes", "database", "schema", "benchmark", "scalab", "sdk"]
    scores["technical"] += sum(3 for t in tech_terms if t in text_lower)

    # Medical signals
    med_terms = ["diagnosis", "prognosis", "pathology", "symptom", "dosage", "patient",
                 "clinical trial", "adverse event", "contraindic", "pharmacol", "therapeutic",
                 "mg/kg", "baseline", "cohort", "randomized", "placebo", "histolog"]
    scores["medical"] += sum(3 for t in med_terms if t in text_lower)

    # Academic signals
    acad_terms = ["abstract", "methodology", "hypothesis", "peer-review", "findings",
                  "literature review", "et al", "citation", "empirical", "p-value",
                  "statistically significant", "longitudinal", "qualitative", "quantitative"]
    scores["academic"] += sum(3 for t in acad_terms if t in text_lower)

    # Financial signals
    fin_terms = ["revenue", "ebitda", "roi", "balance sheet", "cash flow", "fiscal",
                 "quarterly", "shareholder", "dividend", "valuation", "forecast",
                 "profit margin", "equity", "capital expenditure", "net income"]
    scores["financial"] += sum(3 for t in fin_terms if t in text_lower)

    # Marketing signals
    mkt_terms = ["brand", "campaign", "target audience", "engagement", "conversion",
                 "call to action", "cta", "roi", "impression", "click-through",
                 "customer journey", "persona", "content strategy", "social media"]
    scores["marketing"] += sum(3 for t in mkt_terms if t in text_lower)

    best_style = max(scores, key=scores.get)
    return best_style if scores[best_style] >= 6 else "casual"


def detect_user_intent(query: str) -> str:
    """Classify the user's intent from their chat message."""
    q = query.lower().strip()

    rewrite_patterns = ["rewrite", "rephrase", "improve", "edit", "polish", "revise",
                        "make it better", "strengthen", "refine", "rework", "redraft"]
    if any(p in q for p in rewrite_patterns):
        return "rewrite"

    summary_patterns = ["summarize", "summary", "tldr", "tl;dr", "brief", "overview",
                        "key points", "main points", "gist", "digest", "in short"]
    if any(p in q for p in summary_patterns):
        return "summarize"

    compare_patterns = ["compare", "contrast", "difference between", "vs", "versus",
                        "how does .* differ", "similarities"]
    if any(p in q for p in compare_patterns):
        return "compare"

    extract_patterns = ["extract", "pull out", "find all", "list all", "what are the",
                        "give me the", "show me the", "identify all"]
    if any(p in q for p in extract_patterns):
        return "extract"

    analysis_patterns = ["analyze", "analyse", "evaluate", "assess", "review",
                         "examine", "break down", "deep dive", "critique", "audit"]
    if any(p in q for p in analysis_patterns):
        return "analysis"

    # Default: if the query is short / conversational, treat as casual; otherwise analysis
    if len(q.split()) <= 8:
        return "casual"
    return "analysis"


def build_tone_directive(doc_style: str, user_intent: str) -> str:
    """Build a tone adaptation directive to inject into the system prompt."""
    profile = DOCUMENT_STYLE_PROFILES.get(doc_style, DOCUMENT_STYLE_PROFILES["casual"])
    intent_mod = INTENT_TONE_MODIFIERS.get(user_intent, INTENT_TONE_MODIFIERS["casual"])

    return (
        f"TONE ADAPTATION (auto-detected):\n"
        f"- Document Style: {profile['label']}\n"
        f"- Required Tone: {profile['tone']}\n"
        f"- Style Directives: {profile['directives']}\n"
        f"- User Intent: {user_intent.title()}\n"
        f"- Intent Directives: {intent_mod}\n"
        f"\nIMPORTANT: Adapt your response style to match the above tone profile. "
        f"Do NOT use a generic voice — calibrate your vocabulary, structure, and formality "
        f"to feel native to this document domain and the user's current goal.\n"
    )

# Initialize LLM and RAG Components
llm = ChatGoogleGenerativeAI(model="gemini-3.1-flash-lite", temperature=0)

CHROMA_DB_DIR = "chroma_db"
os.makedirs(CHROMA_DB_DIR, exist_ok=True)
embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")
vectorstore = Chroma(persist_directory=CHROMA_DB_DIR, embedding_function=embeddings)
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

class Message(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    messages: List[Message]
    file_id: Optional[str] = None
    conversation_history: Optional[List[Message]] = None

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    file_location = f"{UPLOAD_DIR}/{file.filename}"
    with open(file_location, "wb+") as file_object:
        shutil.copyfileobj(file.file, file_object)
    
    # Process document
    try:
        ext = file.filename.lower().split('.')[-1]
        
        extracted_text = ""
        content_map = {}
        file_id = uuid.uuid4().hex
        
        if ext == "pdf":
            import fitz
            doc = fitz.open(file_location)
            pages = []
            extracted_text_parts = []
            for page_idx, page in enumerate(doc):
                blocks = page.get_text("blocks")
                block_list = []
                for b in blocks:
                    block_list.append({
                        "text": b[4],
                        "bbox": [b[0], b[1], b[2], b[3]],
                        "block_number": b[5],
                        "block_type": b[6]
                    })
                pages.append({
                    "page_index": page_idx,
                    "width": page.rect.width,
                    "height": page.rect.height,
                    "blocks": block_list
                })
                extracted_text_parts.append(page.get_text())
            extracted_text = "\n".join(extracted_text_parts)
            content_map = {"pages": pages}
            
        elif ext == "docx":
            import docx
            doc = docx.Document(file_location)
            paragraphs = []
            extracted_text_parts = []
            for idx, para in enumerate(doc.paragraphs):
                runs_list = []
                for run in para.runs:
                    color_val = None
                    if run.font and run.font.color:
                        try:
                            if run.font.color.rgb:
                                color_val = str(run.font.color.rgb)
                        except Exception:
                            pass
                    
                    runs_list.append({
                        "text": run.text,
                        "bold": run.bold,
                        "italic": run.italic,
                        "font_name": run.font.name if run.font else None,
                        "font_size": run.font.size.pt if (run.font and run.font.size) else None,
                        "color": color_val
                    })
                paragraphs.append({
                    "index": idx,
                    "text": para.text,
                    "runs": runs_list
                })
                extracted_text_parts.append(para.text)
            extracted_text = "\n".join(extracted_text_parts)
            content_map = {"paragraphs": paragraphs}
            
        elif ext == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_location, data_only=True)
            sheets = []
            extracted_text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                cells_list = []
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            val_str = str(cell.value)
                            extracted_text_parts.append(f"Cell {cell.coordinate}: {val_str}")
                            
                            formatting = {}
                            if cell.font:
                                color_val = None
                                if cell.font.color and hasattr(cell.font.color, 'rgb'):
                                    color_val = str(cell.font.color.rgb)
                                formatting = {
                                    "bold": cell.font.bold,
                                    "italic": cell.font.italic,
                                    "font_name": cell.font.name,
                                    "font_size": cell.font.size,
                                    "color": color_val
                                }
                            if cell.number_format:
                                formatting["number_format"] = cell.number_format
                            
                            cells_list.append({
                                "address": cell.coordinate,
                                "value": val_str,
                                "formatting": formatting
                            })
                sheets.append({
                    "sheet_name": sheet_name,
                    "cells": cells_list
                })
            extracted_text = "\n".join(extracted_text_parts)
            content_map = {"sheets": sheets}
            
        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(file_location)
            slides = []
            extracted_text_parts = []
            
            for slide_idx, slide in enumerate(prs.slides):
                shapes_list = []
                
                def crawl_shape(shape):
                    shape_data = {
                        "name": shape.name,
                        "type": str(shape.shape_type),
                        "text": None,
                        "formatting": {}
                    }
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            shape_data["text"] = text
                            extracted_text_parts.append(text)
                        
                        if shape.text_frame.paragraphs:
                            p = shape.text_frame.paragraphs[0]
                            if p.runs:
                                r = p.runs[0]
                                color_val = None
                                if r.font and r.font.color:
                                    try:
                                        if r.font.color.rgb:
                                            color_val = str(r.font.color.rgb)
                                    except Exception:
                                        pass
                                shape_data["formatting"] = {
                                    "bold": r.font.bold if r.font else None,
                                    "italic": r.font.italic if r.font else None,
                                    "font_name": r.font.name if r.font else None,
                                    "font_size": r.font.size.pt if (r.font and r.font.size) else None,
                                    "color": color_val
                                }
                    if shape.has_table:
                        table_cells = []
                        for row in shape.table.rows:
                            row_cells = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                row_cells.append(cell_text)
                                if cell_text:
                                    extracted_text_parts.append(cell_text)
                            table_cells.append(row_cells)
                        shape_data["text"] = json.dumps(table_cells)
                    
                    shapes_list.append(shape_data)
                    
                    if shape.shape_type == 6: # GROUP
                        for subshape in shape.shapes:
                            crawl_shape(subshape)

                for shape in slide.shapes:
                    crawl_shape(shape)
                
                slides.append({
                    "slide_index": slide_idx,
                    "shapes": shapes_list
                })
                
            extracted_text = "\n".join(extracted_text_parts)
            content_map = {"slides": slides}
        else:
            raise HTTPException(
                status_code=400,
                detail="Unsupported file type. Only docx, pdf, pptx, and xlsx files are accepted."
            )
        
        # Save content_map globally keyed by file_id
        content_maps[file_id] = content_map
        
        # Classify document style for tone adaptation
        detected_style = classify_document_style(extracted_text)
        style_label = DOCUMENT_STYLE_PROFILES.get(detected_style, DOCUMENT_STYLE_PROFILES["casual"])["label"]
        
        # Store full text and metadata in memory for initial insights
        uploaded_documents.append({
            "file_id": file_id,
            "filename": file.filename,
            "content": extracted_text,
            "style": detected_style,
            "content_map": content_map,
            "file_path": file_location
        })
        
        # Split into chunks and insert into Vector DB for RAG Chat
        chunks = text_splitter.split_text(extracted_text)
        docs_to_insert = [Document(page_content=chunk, metadata={"filename": file.filename}) for chunk in chunks]
        if docs_to_insert:
            vectorstore.add_documents(docs_to_insert)
        
        # Generate structured summary automatically
        summary_prompt = (
            "You are an expert document analyst. Please analyze the following document. "
            "Identify 1-3 specific sentences or phrases in the text that are weak, unclear, wordy, passive, grammatically incorrect, or low-quality and would benefit from proactive improvement.\n\n"
            "Return a raw JSON object (without markdown code blocks) with EXACTLY these four keys:\n"
            '  "summary": "A concise 2-3 sentence summary.",\n'
            '  "risks": "Any major risks, gaps, or contradictions you detect.",\n'
            '  "recommendations": "A brief expert recommendation based on the content.",\n'
            '  "proactive_suggestions": [\n'
            '    {\n'
            '      "original_text": "The exact weak/unclear text from the document (must match character-for-character, including spacing, to allow search-and-replace).",\n'
            '      "suggested_replacement": "The polished, high-quality replacement text. Must NOT be wrapped in markdown code blocks or quotes.",\n'
            '      "explanation": "Why the original was weak and what was improved."\n'
            '    }\n'
            '  ]\n\n'
            f"Document text:\n{extracted_text[:40000]}"
        )
        try:
            summary_response = llm.invoke([HumanMessage(content=summary_prompt)])
            
            try:
                # Clean potential markdown formatting
                raw_content = summary_response.content.strip()
                if raw_content.startswith("```json"):
                    raw_content = raw_content[7:]
                if raw_content.startswith("```"):
                    raw_content = raw_content[3:]
                if raw_content.endswith("```"):
                    raw_content = raw_content[:-3]
                
                insights = json.loads(raw_content.strip())
                
                # Sanitize suggestions
                if "proactive_suggestions" not in insights:
                    insights["proactive_suggestions"] = []
                else:
                    for sug in insights["proactive_suggestions"]:
                        orig = sug.get("original_text", "")
                        sug["original_text"] = orig.strip() if isinstance(orig, str) else ""
                        rep = sug.get("suggested_replacement", "")
                        if isinstance(rep, str):
                            rep = rep.strip()
                            if rep.startswith("```"):
                                rep = re.sub(r"^```[a-zA-Z]*\n?", "", rep)
                            if rep.endswith("```"):
                                rep = re.sub(r"\n?```$", "", rep)
                            sug["suggested_replacement"] = rep
                        else:
                            sug["suggested_replacement"] = ""
            except Exception as e:
                # Fallback if JSON parsing fails
                insights = {
                    "summary": summary_response.content,
                    "risks": "Could not structure risks. See summary.",
                    "recommendations": "Could not structure recommendations. See summary.",
                    "proactive_suggestions": []
                }
        except Exception as err:
            print(f"Summary generation failed (likely rate-limited): {err}")
            insights = {
                "summary": "Summary generation was skipped or rate-limited.",
                "risks": "N/A",
                "recommendations": "N/A",
                "proactive_suggestions": []
            }
        
        # Calculate structured counts for the short summary
        word_count = len(extracted_text.split())
        short_summary = {
            "file_type": ext,
            "word_count": word_count
        }
        if ext == "pdf":
            short_summary["page_count"] = len(content_map.get("pages", []))
        elif ext == "docx":
            short_summary["paragraph_count"] = len(content_map.get("paragraphs", []))
        elif ext == "xlsx":
            short_summary["sheet_count"] = len(content_map.get("sheets", []))
        elif ext == "pptx":
            short_summary["slide_count"] = len(content_map.get("slides", []))

        return {
            "message": f"Successfully uploaded {file.filename}", 
            "file_id": file_id,
            "short_summary": short_summary,
            "insights": insights, 
            "filename": file.filename,
            "url": f"http://localhost:8000/files/{file.filename}",
            "content": extracted_text,
            "content_map": content_map,
            "document_style": detected_style,
            "document_style_label": style_label
        }
    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.get("/content-map/{file_id}")
async def get_content_map(file_id: str):
    """Retrieve a stored content map by file_id."""
    if file_id not in content_maps:
        raise HTTPException(status_code=404, detail="Content map not found for this file_id.")
    
    # Find metadata from uploaded_documents
    file_type = None
    filename = None
    for doc in uploaded_documents:
        if doc.get("file_id") == file_id:
            filename = doc.get("filename", "")
            ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
            file_type = ext
            break
    
    return {
        "file_id": file_id,
        "file_type": file_type,
        "filename": filename,
        "content_map": content_maps[file_id]
    }

STOP_WORDS = {
    "the", "a", "an", "and", "or", "but", "if", "then", "else", "of", "at", "by", 
    "for", "with", "about", "against", "between", "into", "through", "during", 
    "before", "after", "above", "below", "to", "from", "up", "down", "in", "out", 
    "on", "off", "over", "under", "again", "further", "then", "once", "here", 
    "there", "when", "where", "why", "how", "all", "any", "both", "each", "few", 
    "more", "most", "other", "some", "such", "no", "nor", "not", "only", "own", 
    "same", "so", "than", "too", "very", "s", "t", "can", "will", "just", "don", 
    "should", "now", "i", "me", "my", "myself", "we", "our", "ours", "ourselves", 
    "you", "your", "yours", "yourself", "yourselves", "he", "him", "his", "himself", 
    "she", "her", "hers", "herself", "it", "its", "itself", "they", "them", "their", 
    "theirs", "themselves", "what", "which", "who", "whom", "this", "that", "these", 
    "those", "am", "is", "are", "was", "were", "be", "been", "being", "have", "has", 
    "had", "having", "do", "does", "did", "doing", "would", "should", "could", "ought"
}

def flatten_content_map(content_map: dict) -> List[Dict[str, str]]:
    items = []
    if not content_map or not isinstance(content_map, dict):
        return items
        
    # PDF
    if "pages" in content_map:
        for page in content_map["pages"]:
            page_num = page.get("page_index", 0) + 1
            for block in page.get("blocks", []):
                text = block.get("text", "").strip()
                if text:
                    items.append({
                        "text": text,
                        "reference": f"Page {page_num}, Block {block.get('block_number', 0)}"
                    })
                    
    # DOCX or text submissions
    elif "paragraphs" in content_map:
        for para in content_map["paragraphs"]:
            text = para.get("text", "").strip()
            if text:
                items.append({
                    "text": text,
                    "reference": f"Paragraph {para.get('index', 0) + 1}"
                })
                
    # PPTX
    elif "slides" in content_map:
        for slide in content_map["slides"]:
            slide_num = slide.get("slide_index", 0) + 1
            for shape in slide.get("shapes", []):
                text = shape.get("text", "")
                if text:
                    text = text.strip()
                    # Check if text is a JSON table (from PPTX)
                    if text.startswith("[[") and text.endswith("]]"):
                        try:
                            rows = json.loads(text)
                            table_lines = []
                            for row in rows:
                                table_lines.append(" | ".join(str(cell) for cell in row))
                            text = "\n".join(table_lines)
                        except Exception:
                            pass
                    name = shape.get("name", "Shape")
                    items.append({
                        "text": text,
                        "reference": f"Slide {slide_num}, {name}"
                    })
                    
    # XLSX
    elif "sheets" in content_map:
        for sheet in content_map["sheets"]:
            sheet_name = sheet.get("sheet_name", "Sheet")
            for cell in sheet.get("cells", []):
                val = cell.get("value", "").strip()
                if val:
                    items.append({
                        "text": val,
                        "reference": f"Sheet '{sheet_name}', Cell {cell.get('address', '')}"
                    })
                    
    return items

def tokenize_text(text: str) -> List[str]:
    words = re.findall(r'[a-zA-Z0-9]+', text.lower())
    return [w for w in words if len(w) >= 2 and w not in STOP_WORDS]

def build_context_from_content_map(file_id: str, query: str) -> str:
    if file_id not in content_maps:
        return ""
    
    content_map = content_maps[file_id]
    items = flatten_content_map(content_map)
    if not items:
        return ""
        
    # Count total words
    total_words = sum(len(item["text"].split()) for item in items)
    
    selected_items = []
    
    if total_words < 6000:
        # Include everything
        selected_items = items
    else:
        # Keyword relevance selection
        query_tokens = tokenize_text(query)
        if not query_tokens:
            words = re.findall(r'[a-zA-Z0-9]+', query.lower())
            query_tokens = [w for w in words if len(w) >= 2]
            
        for item in items:
            text_words = re.findall(r'[a-zA-Z0-9]+', item["text"].lower())
            score = 0
            for token in query_tokens:
                score += text_words.count(token)
            item["score"] = score
            
        # Check if all scores are 0 (no keyword overlap)
        all_zero = all(item["score"] == 0 for item in items)
        
        if all_zero:
            # Fallback to sequential selection from the beginning
            selected_items = []
            curr_words = 0
            for item in items:
                item_word_cnt = len(item["text"].split())
                if curr_words + item_word_cnt > 6000:
                    break
                selected_items.append(item)
                curr_words += item_word_cnt
        else:
            # Sort by score descending, keeping original order as secondary key to stabilize sorting
            for idx, item in enumerate(items):
                item["original_index"] = idx
                
            sorted_items = sorted(items, key=lambda x: (-x["score"], x["original_index"]))
            
            # Select top items up to 6000 words
            selected_items = []
            curr_words = 0
            for item in sorted_items:
                item_word_cnt = len(item["text"].split())
                if curr_words + item_word_cnt > 6000 and len(selected_items) > 0:
                    break
                selected_items.append(item)
                curr_words += item_word_cnt
                
            # Re-sort selected items by original index to keep document flow/chronology
            selected_items = sorted(selected_items, key=lambda x: x["original_index"])
            
    # Format the selected items
    formatted_parts = []
    for item in selected_items:
        formatted_parts.append(
            f"--- SECTION: {item['reference']} ---\n"
            f"{item['text']}\n"
            f"--- END SECTION: {item['reference']} ---"
        )
    return "\n\n".join(formatted_parts)

def condense_query(query: str, history: List[Message]) -> str:
    if not history:
        return query
        
    history_str = ""
    for msg in history:
        role_label = "User" if msg.role == "user" else "Assistant"
        history_str += f"{role_label}: {msg.content}\n"
        
    condense_prompt = (
        "Given the following conversation history and a follow-up question, "
        "rephrase the follow-up question to be a standalone question (containing all necessary context from the conversation history, including pronouns and subject references).\n"
        "Return ONLY the standalone question, with no introduction or preamble.\n\n"
        f"Conversation History:\n{history_str}\n"
        f"Follow-up Question: {query}\n"
        "Standalone Question:"
    )
    
    try:
        response = llm.invoke([HumanMessage(content=condense_prompt)])
        condensed = response.content.strip()
        if condensed:
            return condensed
    except Exception:
        pass
        
    return query


def condense_instruction(instruction: str, history: List[Message]) -> str:
    if not history:
        return instruction
        
    history_str = ""
    for msg in history:
        role_label = "User" if msg.role == "user" else "Assistant"
        history_str += f"{role_label}: {msg.content}\n"
        
    condense_prompt = (
        "Given the following conversation history and a follow-up edit instruction, "
        "rephrase the follow-up instruction to be a standalone edit instruction that explicitly names the targeted text, location, or content reference (resolving pronouns like 'it', 'that', 'this name', 'this line' based on the history).\n"
        "Return ONLY the standalone rephrased instruction, with no introduction or preamble.\n\n"
        f"Conversation History:\n{history_str}\n"
        f"Follow-up Edit Instruction: {instruction}\n"
        "Standalone Edit Instruction:"
    )
    
    try:
        response = llm.invoke([HumanMessage(content=condense_prompt)])
        condensed = response.content.strip()
        if condensed:
            return condensed
    except Exception as e:
        print(f"Condensation failed: {e}")
        
    return instruction

@app.post("/chat")
async def chat(request: ChatRequest):
    try:
        if not uploaded_documents:
            return {"answer": "Please upload a document first."}
            
        # Get the latest user query for retrieval
        latest_query = request.messages[-1].content if request.messages else ""
        
        # Check if file_id is provided and valid in content_maps
        is_content_map_query = False
        retrieved_context = ""
        
        search_query = latest_query
        
        # Construct full history for condensation and LLM context
        full_history = []
        if request.conversation_history:
            full_history.extend(request.conversation_history)
        if len(request.messages) > 1:
            full_history.extend(request.messages[:-1])
            
        if request.file_id and request.file_id in content_maps:
            is_content_map_query = True
            
            if full_history:
                search_query = condense_query(latest_query, full_history)
                
            retrieved_context = build_context_from_content_map(request.file_id, search_query)
            
        if is_content_map_query:
            # TODO: Replace with the user-provided system prompt.
            # Below is a temporary placeholder instructing the model to answer ONLY from the provided document content,
            # and say clearly if the answer isn't in the document, rather than guessing.
            system_prompt = (
                "You are an AI Document Assistant.\n"
                "CRITICAL INSTRUCTIONS:\n"
                "1. Answer the user's question ONLY based on the provided document content below.\n"
                "2. Do NOT guess, hypothesize, extrapolate, or use any general or external knowledge.\n"
                "3. If the answer is not contained in the provided document content, you MUST state clearly: "
                "\"I cannot find the answer to this question in the provided document content.\"\n"
                "4. Cite the source reference inline when referencing facts (e.g. 'The launch date is December 25 [Paragraph 2]').\n"
                "5. Identify which sections (using their exact references, e.g. 'Paragraph 1', 'Page 2, Block 3') you used to formulate your answer.\n\n"
                "You MUST return your response as a raw JSON object (do not wrap in markdown code blocks or any other text) with the following structure:\n"
                "{\n"
                "  \"answer\": \"The text of your answer, citing sources inline where appropriate, or the fallback message if not found.\",\n"
                "  \"sources\": [\"Reference of section 1\", \"Reference of section 2\"]\n"
                "}\n\n"
                "Provided Document Content:\n"
                f"{retrieved_context}"
            )
        else:
            # Detect user intent and determine dominant document style
            user_intent = detect_user_intent(latest_query)
            
            # Determine document style from uploaded documents
            # Use the most recently uploaded document's style, or aggregate
            doc_style = "casual"
            if uploaded_documents:
                style_counts: Dict[str, int] = {}
                for doc in uploaded_documents:
                    s = doc.get("style", "casual")
                    style_counts[s] = style_counts.get(s, 0) + 1
                doc_style = max(style_counts, key=style_counts.get)
            
            # Build adaptive tone directive
            tone_directive = build_tone_directive(doc_style, user_intent)
            
            # Perform Vector Search
            retrieved_docs = vectorstore.similarity_search(latest_query, k=5)
            
            # Construct context from retrieved chunks only
            context_parts = []
            for i, doc in enumerate(retrieved_docs):
                context_parts.append(f"--- CHUNK {i+1} [Source: {doc.metadata.get('filename', 'Unknown')}] ---\n{doc.page_content}\n--- END CHUNK {i+1} ---")
                
            retrieved_context = "\n\n".join(context_parts)
            
            system_prompt = (
                "You are an elite AI Document Intelligence and Editing Assistant — a world-class analyst, consultant, and professional editor rolled into one. "
                "Your capabilities span deep document analysis, strategic consulting, and masterful text editing.\n\n"
                
                f"{tone_directive}\n\n"
                
                "DOCUMENT INTELLIGENCE:\n"
                "- Read, understand, and synthesize information strictly from the provided retrieved document chunks.\n"
                "- Always prioritize the retrieved chunks over your general knowledge. Never hallucinate data not present in them.\n"
                "- Cross-reference information across multiple chunks to identify patterns, connections, and discrepancies.\n"
                "- Proactively detect risks, gaps, contradictions, and hidden opportunities within the content.\n"
                "- Generate highly actionable insights backed by specific evidence from the documents.\n\n"
                
                "ANALYTICAL REASONING:\n"
                "- Use rigorous step-by-step reasoning before answering complex questions.\n"
                "- Strictly prioritize accuracy over assumptions. State clearly when information is insufficient.\n"
                "- Provide expert recommendations with clear rationale and supporting evidence.\n"
                "- Compare and contrast information across documents when multiple sources are available.\n\n"
                
                "CLARIFICATION & DIRECTIVES:\n"
                "- If the user's query, task, or editing instruction is unclear, ambiguous, or contains contradictory commands, do NOT guess. Ask a direct, polite clarifying question to resolve the ambiguity.\n\n"
                
                "EDITING & WRITING:\n"
                "- You are fluent across all content domains: legal, academic, business, technical, medical, financial, marketing, and creative. Automatically detect the domain and apply the appropriate conventions, terminology, and standards.\n"
                "- THINK BEFORE EDITING. Before making any changes, deeply understand the intent, purpose, audience, and desired outcome of the text. Never edit reflexively.\n"
                "- When asked to write, rewrite, or edit a section of text, return the final polished replacement text ONLY. Do NOT include any conversational preamble, introduction, explanation, notes, or postscript. The response must contain only the direct copy-pasteable replacement content.\n"
                "- Preserve the original intent and meaning while making it significantly stronger.\n"
                "- Automatically chooses the most logical and effective wording — do not present multiple options; deliver the single best version decisively.\n"
                "- NEVER paraphrase blindly. Every change must upgrade the meaning, clarity, and impact of the original text.\n"
                "- Ensure replacements fit perfectly into the surrounding document context — match the style, terminology, and structure of the rest of the document so the edit feels seamless, not bolted on.\n"
                "- If the surrounding context is weak or inconsistent, proactively rewrite it as well to improve overall coherence and flow.\n"
                "- Produce professional summaries, executive briefs, and polished reports on demand.\n\n"
                
                "If the answer to the user's question is not contained within the provided retrieved chunks, explicitly state that you cannot answer based on the uploaded documents.\n\n"
                "Here are the retrieved chunks most relevant to the user's query:\n\n"
                f"{retrieved_context}"
            )
        
        langchain_messages = [SystemMessage(content=system_prompt)]
        
        # Prepend conversation_history if provided
        if request.conversation_history:
            for msg in request.conversation_history:
                if msg.role == "user":
                    langchain_messages.append(HumanMessage(content=msg.content))
                elif msg.role == "bot":
                    langchain_messages.append(AIMessage(content=msg.content))
                    
        for msg in request.messages:
            if msg.role == "user":
                langchain_messages.append(HumanMessage(content=msg.content))
            elif msg.role == "bot":
                langchain_messages.append(AIMessage(content=msg.content))
        
        response = llm.invoke(langchain_messages)
        answer_text = response.content.strip()
        sources_list = []
        
        if is_content_map_query:
            try:
                raw_content = answer_text
                if raw_content.startswith("```json"):
                    raw_content = raw_content[7:]
                if raw_content.startswith("```"):
                    raw_content = raw_content[3:]
                if raw_content.endswith("```"):
                    raw_content = raw_content[:-3]
                
                res_data = json.loads(raw_content.strip())
                answer_text = res_data.get("answer", answer_text)
                sources_list = res_data.get("sources", [])
                if not isinstance(sources_list, list):
                    sources_list = [sources_list] if sources_list else []
            except Exception:
                refs = re.findall(r'\[(Paragraph \d+|Page \d+, Block \d+|Slide \d+, [^\]]+|Sheet \'[^\']+\', Cell [A-Z]+\d+)\]', answer_text)
                if refs:
                    sources_list = list(set(refs))
        else:
            sources_list = list(set(doc.metadata.get('filename', 'Unknown') for doc in retrieved_docs))
            
        return {"answer": answer_text, "sources": sources_list}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


class TextSubmission(BaseModel):
    title: str
    text: str

@app.post("/submit-text")
async def submit_text(submission: TextSubmission):
    """Accept pasted text (partial, full document, or selected section) for analysis."""
    try:
        extracted_text = submission.text.strip()
        if not extracted_text:
            raise HTTPException(status_code=400, detail="Text content is empty.")
        
        title = submission.title.strip() or "Pasted Text"
        
        # Classify document style for tone adaptation
        detected_style = classify_document_style(extracted_text)
        style_label = DOCUMENT_STYLE_PROFILES.get(detected_style, DOCUMENT_STYLE_PROFILES["casual"])["label"]
        
        file_id = uuid.uuid4().hex
        paragraphs_list = []
        for idx, line in enumerate(extracted_text.split("\n")):
            if line.strip():
                paragraphs_list.append({
                    "index": idx,
                    "text": line,
                    "runs": []
                })
        content_map = {"paragraphs": paragraphs_list}
        content_maps[file_id] = content_map

        file_location = f"{UPLOAD_DIR}/{file_id}.txt"
        with open(file_location, "w", encoding="utf-8") as f:
            f.write(extracted_text)

        # Store in memory
        uploaded_documents.append({
            "file_id": file_id,
            "filename": title,
            "content": extracted_text,
            "style": detected_style,
            "content_map": content_map,
            "file_path": file_location
        })
        
        # Split into chunks and insert into Vector DB for RAG Chat
        chunks = text_splitter.split_text(extracted_text)
        docs_to_insert = [Document(page_content=chunk, metadata={"filename": title}) for chunk in chunks]
        if docs_to_insert:
            vectorstore.add_documents(docs_to_insert)
        
        # Generate structured summary
        summary_prompt = (
            "You are an expert document analyst. Please analyze the following text. "
            "Identify 1-3 specific sentences or phrases in the text that are weak, unclear, wordy, passive, grammatically incorrect, or low-quality and would benefit from proactive improvement.\n\n"
            "Return a raw JSON object (without markdown code blocks) with EXACTLY these four keys:\n"
            '  "summary": "A concise 2-3 sentence summary.",\n'
            '  "risks": "Any major risks, gaps, or contradictions you detect.",\n'
            '  "recommendations": "A brief expert recommendation based on the content.",\n'
            '  "proactive_suggestions": [\n'
            '    {\n'
            '      "original_text": "The exact weak/unclear text from the document (must match character-for-character, including spacing, to allow search-and-replace).",\n'
            '      "suggested_replacement": "The polished, high-quality replacement text. Must NOT be wrapped in markdown code blocks or quotes.",\n'
            '      "explanation": "Why the original was weak and what was improved."\n'
            '    }\n'
            '  ]\n\n'
            f"Text:\n{extracted_text[:40000]}"
        )
        summary_response = llm.invoke([HumanMessage(content=summary_prompt)])
        
        try:
            raw_content = summary_response.content.strip()
            if raw_content.startswith("```json"):
                raw_content = raw_content[7:]
            if raw_content.startswith("```"):
                raw_content = raw_content[3:]
            if raw_content.endswith("```"):
                raw_content = raw_content[:-3]
            insights = json.loads(raw_content.strip())
            
            # Sanitize suggestions
            if "proactive_suggestions" not in insights:
                insights["proactive_suggestions"] = []
            else:
                for sug in insights["proactive_suggestions"]:
                    orig = sug.get("original_text", "")
                    sug["original_text"] = orig.strip() if isinstance(orig, str) else ""
                    rep = sug.get("suggested_replacement", "")
                    if isinstance(rep, str):
                        rep = rep.strip()
                        if rep.startswith("```"):
                            rep = re.sub(r"^```[a-zA-Z]*\n?", "", rep)
                        if rep.endswith("```"):
                            rep = re.sub(r"\n?```$", "", rep)
                        sug["suggested_replacement"] = rep
                    else:
                        sug["suggested_replacement"] = ""
        except Exception:
            insights = {
                "summary": summary_response.content,
                "risks": "Could not structure risks. See summary.",
                "recommendations": "Could not structure recommendations. See summary.",
                "proactive_suggestions": []
            }
        
        return {
            "message": f"Successfully processed: {title}",
            "file_id": file_id,
            "title": title,
            "content": extracted_text,
            "content_map": content_map,
            "insights": insights,
            "document_style": detected_style,
            "document_style_label": style_label
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing text: {str(e)}")


class EditRequest(BaseModel):
    selected_text: str
    full_context: str = ""
    instruction: str = "improve"

@app.post("/edit-section")
async def edit_section(request: EditRequest):
    """AI-powered editing of a selected text section."""
    try:
        if not request.selected_text.strip():
            raise HTTPException(status_code=400, detail="No text selected for editing.")
        
        # Detect the style of the surrounding context for tone-matched editing
        context_style = classify_document_style(request.full_context or request.selected_text)
        style_profile = DOCUMENT_STYLE_PROFILES.get(context_style, DOCUMENT_STYLE_PROFILES["casual"])
        
        edit_prompt = (
            f"You are an elite professional editor specializing in {style_profile['label']} documents.\n\n"
            f"TONE: {style_profile['tone']}\n"
            f"STYLE RULES: {style_profile['directives']}\n\n"
        )
        
        if request.full_context:
            edit_prompt += f"SURROUNDING DOCUMENT CONTEXT (for style matching):\n{request.full_context[:5000]}\n\n"
        
        edit_prompt += (
            f"SELECTED TEXT TO EDIT:\n\"{request.selected_text}\"\n\n"
            f"INSTRUCTION FROM USER:\n\"{request.instruction}\"\n\n"
            "Analyze the user's instruction. If the instruction is unclear, ambiguous, too vague to perform the edit, or contains contradictory commands, set \"needs_clarification\" to true and write a clarifying question under \"clarification_question\" asking the user to explain their intent.\n"
            "Otherwise, set \"needs_clarification\" to false and provide the edited text.\n\n"
            "Return a raw JSON object (without markdown code blocks) with EXACTLY these four keys:\n"
            '  "edited_text": "The final polished replacement text. Must be a direct drop-in replacement. Set to empty string if needs_clarification is true.",\n'
            '  "explanation": "A brief 1-sentence explanation of what was improved (or why clarification is needed).",\n'
            '  "needs_clarification": true/false,\n'
            '  "clarification_question": "If needs_clarification is true, specify a clarifying question asking the user for details. Otherwise, set this to empty string."\n\n'
            "Rules for the edited text:\n"
            "- Return the final polished replacement text ONLY. Do NOT include any introduction, conversational preamble, explanation, or notes in the edited_text value.\n"
            "- Do NOT wrap the edited_text value in markdown code blocks (like ```) or quotes. The text must be returned as plain text inside the JSON string, completely clean and ready to paste directly into the document.\n"
            "- Preserve the original meaning and intent unless the user explicitly instructs otherwise.\n"
            "- Match the tone and style of the surrounding document.\n"
            "- Fix grammar, strengthen phrasing, improve clarity.\n"
            "- Return ONLY the replacement text for the selected section, not the entire document.\n"
        )
        
        response = llm.invoke([HumanMessage(content=edit_prompt)])
        
        try:
            raw_content = response.content.strip()
            if raw_content.startswith("```json"):
                raw_content = raw_content[7:]
            if raw_content.startswith("```"):
                raw_content = raw_content[3:]
            if raw_content.endswith("```"):
                raw_content = raw_content[:-3]
            result = json.loads(raw_content.strip())
            
            edited_text = result.get("edited_text", "")
            if isinstance(edited_text, str):
                edited_text = edited_text.strip()
                if edited_text.startswith("```"):
                    edited_text = re.sub(r"^```[a-zA-Z]*\n?", "", edited_text)
                if edited_text.endswith("```"):
                    edited_text = re.sub(r"\n?```$", "", edited_text)
                    
            return {
                "edited_text": edited_text,
                "explanation": result.get("explanation", ""),
                "needs_clarification": result.get("needs_clarification", False),
                "clarification_question": result.get("clarification_question", "")
            }
        except Exception:
            edited_text = response.content.strip()
            if edited_text.startswith("```"):
                edited_text = re.sub(r"^```[a-zA-Z]*\n?", "", edited_text)
            if edited_text.endswith("```"):
                edited_text = re.sub(r"\n?```$", "", edited_text)
            return {
                "edited_text": edited_text,
                "explanation": "AI editing applied.",
                "needs_clarification": False,
                "clarification_question": ""
            }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error editing section: {str(e)}")


class LocateRequest(BaseModel):
    file_id: str
    instruction: str
    conversation_history: Optional[List[Message]] = None

@app.post("/edit/locate")
async def edit_locate(request: LocateRequest):
    """Identify which sections are targeted by a natural-language edit instruction."""
    try:
        file_id = request.file_id
        instruction = request.instruction.strip()
        if request.conversation_history:
            instruction = condense_instruction(instruction, request.conversation_history)
        
        if file_id not in content_maps:
            raise HTTPException(status_code=404, detail="File content map not found.")
            
        content_map = content_maps[file_id]
        items = flatten_content_map(content_map)
        if not items:
            return {
                "matched_locations": [],
                "needs_clarification": False,
                "clarification_question": ""
            }
            
        # Format the items list for the LLM
        items_str_parts = []
        for idx, item in enumerate(items):
            items_str_parts.append(
                f"--- ITEM REFERENCE: {item['reference']} ---\n"
                f"{item['text']}\n"
                f"--- END ITEM REFERENCE: {item['reference']} ---"
            )
        items_str = "\n\n".join(items_str_parts)
        
        system_prompt = (
            "You identify exactly which part of a document an edit instruction refers to, given a structured map of paragraphs/slides/cells/pages.\n"
            "1. Return the single best-matching location with its exact reference.\n"
            "2. If multiple locations plausibly match, don't guess — say it's ambiguous and list candidates with ~10 words of context each.\n"
            "3. If nothing matches, say so rather than inventing a location.\n"
            "4. Don't draft replacement text here — only locate.\n\n"
            "You MUST return your response as a raw JSON object (do not wrap in markdown code blocks or any other text) with the following structure:\n"
            "{\n"
            "  \"matched_locations\": [\n"
            "    {\n"
            "      \"reference\": \"The exact reference string of the item (e.g. 'Paragraph 5', 'Page 1, Block 2')\",\n"
            "      \"text\": \"The exact original text of that item\"\n"
            "    }\n"
            "  ],\n"
            "  \"needs_clarification\": true/false,\n"
            "  \"clarification_question\": \"Your clarifying question here if needs_clarification is true, listing candidates with ~10 words of context each, otherwise empty string.\"\n"
            "}"
        )
        
        human_content = (
            f"Here is the list of document items with their reference names and original text:\n\n"
            f"{items_str[:40000]}\n\n"
            f"USER EDIT INSTRUCTION:\n\"{instruction}\""
        )
        
        response = llm.invoke([
            SystemMessage(content=system_prompt),
            HumanMessage(content=human_content)
        ])
        
        raw_content = response.content.strip()
        
        # Clean markdown formatting if present
        if raw_content.startswith("```json"):
            raw_content = raw_content[7:]
        if raw_content.startswith("```"):
            raw_content = raw_content[3:]
        if raw_content.endswith("```"):
            raw_content = raw_content[:-3]
            
        try:
            result = json.loads(raw_content.strip())
            matched_locations = result.get("matched_locations", [])
            needs_clarification = result.get("needs_clarification", False)
            clarification_question = result.get("clarification_question", "")
            
            validated_locations = []
            if isinstance(matched_locations, list):
                for loc in matched_locations:
                    if isinstance(loc, dict) and "reference" in loc and "text" in loc:
                        validated_locations.append({
                            "reference": loc["reference"],
                            "text": loc["text"]
                        })
            
            return {
                "matched_locations": validated_locations,
                "needs_clarification": bool(needs_clarification),
                "clarification_question": str(clarification_question)
            }
        except Exception as parse_err:
            return {
                "matched_locations": [],
                "needs_clarification": True,
                "clarification_question": f"Could not parse locate results: {raw_content}"
            }
            
    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error locating edit target: {str(e)}")


class DraftRequest(BaseModel):
    file_id: str
    reference: str
    instruction: str
    conversation_history: Optional[List[Message]] = None

@app.post("/edit/draft")
async def edit_draft(request: DraftRequest):
    """Draft a replacement for a located reference in a document based on instruction."""
    try:
        file_id = request.file_id
        reference = request.reference.strip()
        instruction = request.instruction.strip()
        if request.conversation_history:
            instruction = condense_instruction(instruction, request.conversation_history)
        
        if file_id not in content_maps:
            raise HTTPException(status_code=404, detail="File content map not found.")
            
        content_map = content_maps[file_id]
        items = flatten_content_map(content_map)
        
        target_item = None
        for item in items:
            if item["reference"] == reference:
                target_item = item
                break
                
        if not target_item:
            raise HTTPException(status_code=404, detail=f"Target reference '{reference}' not found in document.")
            
        full_context = "\n\n".join(item["text"] for item in items)
        
        system_prompt = (
            "You draft a replacement for one specific, already-located piece of a document, per a user's edit instruction.\n"
            "1. Match the tone, formality, and rhythm of the surrounding document — write as the original author would, not as generic new content.\n"
            "2. Keep length close to the original unless told to expand/shorten significantly.\n"
            "3. Change only what was asked — don't \"improve\" unrelated wording or fix unrelated issues.\n"
            "4. Output ONLY the replacement text — no preamble, no quotes around it, nothing else.\n"
            "5. If the instruction is too vague to draft confidently, ask one clarifying question instead of guessing.\n\n"
            "Special Output Rule:\n"
            "If you need to ask a clarifying question, prefix it with 'CLARIFICATION_REQUIRED: '. Otherwise, output only the raw replacement text."
        )
        
        human_content = (
            f"DOCUMENT SURROUNDING CONTEXT:\n{full_context[:10000]}\n\n"
            f"TARGET SECTION REFERENCE: {reference}\n"
            f"ORIGINAL TEXT OF SECTION:\n\"{target_item['text']}\"\n\n"
            f"EDIT INSTRUCTION:\n\"{instruction}\""
        )
        
        response = llm.invoke([
            SystemMessage(content=system_prompt),
            HumanMessage(content=human_content)
        ])
        
        content = response.content.strip()
        
        if content.startswith("```"):
            lines = content.split("\n")
            if lines[0].startswith("```"):
                content = "\n".join(lines[1:-1]) if lines[-1].startswith("```") else "\n".join(lines[1:])
            content = content.strip()
            
        if content.startswith("CLARIFICATION_REQUIRED:"):
            clar_q = content.replace("CLARIFICATION_REQUIRED:", "").strip()
            return {
                "reference": reference,
                "original_text": target_item["text"],
                "draft_text": "",
                "needs_clarification": True,
                "clarification_question": clar_q
            }
        else:
            return {
                "reference": reference,
                "original_text": target_item["text"],
                "draft_text": content,
                "needs_clarification": False,
                "clarification_question": ""
            }
            
    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating draft: {str(e)}")


class ApplyRequest(BaseModel):
    file_id: str
    location: str
    text: str

@app.post("/edit/apply")
async def edit_apply(request: ApplyRequest):
    """Apply an approved edit text to a real document, keeping original format and saving as new version."""
    try:
        file_id = request.file_id
        location = request.location.strip()
        text = request.text.strip()
        
        doc_meta = None
        for doc in uploaded_documents:
            if doc.get("file_id") == file_id:
                doc_meta = doc
                break
                
        if not doc_meta:
            raise HTTPException(status_code=404, detail="Uploaded document metadata not found.")
            
        file_path = doc_meta.get("file_path")
        if not file_path or not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="Source file not found on disk.")
            
        ext = file_path.lower().split('.')[-1]
        
        if ext == "docx":
            import docx
            doc_obj = docx.Document(file_path)
            
            match = re.match(r"Paragraph\s+(\d+)", location)
            if not match:
                raise HTTPException(status_code=400, detail=f"Invalid location reference for docx: '{location}'")
            para_idx = int(match.group(1)) - 1
            
            if para_idx < 0 or para_idx >= len(doc_obj.paragraphs):
                raise HTTPException(status_code=400, detail=f"Paragraph index {para_idx+1} out of bounds.")
                
            p = doc_obj.paragraphs[para_idx]
            if p.runs:
                p.runs[0].text = text
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = text
                
        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            
            match = re.match(r"Slide\s+(\d+),\s*(.+)", location)
            if not match:
                raise HTTPException(status_code=400, detail=f"Invalid location reference for pptx: '{location}'")
            slide_idx = int(match.group(1)) - 1
            shape_name = match.group(2).strip()
            
            if slide_idx < 0 or slide_idx >= len(prs.slides):
                raise HTTPException(status_code=400, detail=f"Slide index {slide_idx+1} out of bounds.")
                
            slide = prs.slides[slide_idx]
            
            target_shape = None
            for shape in slide.shapes:
                if shape.name == shape_name:
                    target_shape = shape
                    break
                    
            if not target_shape:
                raise HTTPException(status_code=400, detail=f"Shape '{shape_name}' not found on Slide {slide_idx+1}")
                
            if target_shape.has_text_frame:
                tf = target_shape.text_frame
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    if p.runs:
                        p.runs[0].text = text
                        for r in p.runs[1:]:
                            r.text = ""
                        for extra_p in tf.paragraphs[1:]:
                            extra_p.text = ""
                    else:
                        tf.text = text
                else:
                    tf.text = text
            else:
                raise HTTPException(status_code=400, detail=f"Shape '{shape_name}' does not support text.")
                
        elif ext == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            
            match = re.match(r"Sheet\s+'([^']+)',\s*Cell\s*([A-Z]+\d+)", location)
            if not match:
                raise HTTPException(status_code=400, detail=f"Invalid location reference for xlsx: '{location}'")
            sheet_name = match.group(1)
            cell_address = match.group(2)
            
            if sheet_name not in wb.sheetnames:
                raise HTTPException(status_code=400, detail=f"Sheet '{sheet_name}' not found in workbook.")
                
            ws = wb[sheet_name]
            try:
                if "." in text:
                    ws[cell_address].value = float(text)
                else:
                    ws[cell_address].value = int(text)
            except ValueError:
                ws[cell_address].value = text
                
        elif ext == "pdf":
            import fitz
            doc_obj = fitz.open(file_path)
            
            match = re.match(r"Page\s+(\d+),\s*Block\s+(\d+)", location)
            if not match:
                raise HTTPException(status_code=400, detail=f"Invalid location reference for pdf: '{location}'")
            page_num = int(match.group(1)) - 1
            block_num = int(match.group(2))
            
            if page_num < 0 or page_num >= len(doc_obj):
                raise HTTPException(status_code=400, detail=f"Page number {page_num+1} out of bounds.")
                
            page = doc_obj[page_num]
            blocks_dict = page.get_text("dict")["blocks"]
            target_block = None
            
            if block_num >= 0 and block_num < len(blocks_dict):
                target_block = blocks_dict[block_num]
                
            if not target_block or target_block.get("type") != 0:
                raise HTTPException(status_code=400, detail=f"Text block {block_num} not found on Page {page_num+1}.")
                
            bbox = target_block["bbox"]
            rect = fitz.Rect(bbox)
            
            font_size = 11
            font_name = "helv"
            font_color = (0, 0, 0)
            
            if target_block.get("lines"):
                line = target_block["lines"][0]
                if line.get("spans"):
                    span = line["spans"][0]
                    font_size = span.get("size", font_size)
                    raw_font = span.get("font", "").lower()
                    if "times" in raw_font:
                        font_name = "tiro" if "italic" in raw_font else "times"
                    elif "courier" in raw_font or "mono" in raw_font:
                        font_name = "couri"
                    elif "bold" in raw_font:
                        font_name = "hebo"
                    else:
                        font_name = "helv"
                        
                    col_int = span.get("color", 0)
                    r = ((col_int >> 16) & 255) / 255.0
                    g = ((col_int >> 8) & 255) / 255.0
                    b = (col_int & 255) / 255.0
                    font_color = (r, g, b)
            
            # Apply redactions to scrub and overlay new text
            page.add_redact_annot(rect, text=text, fontname=font_name, fontsize=font_size, text_color=font_color, fill=(1, 1, 1))
            page.apply_redactions()
            
        elif ext == "txt":
            match = re.match(r"Paragraph\s+(\d+)", location)
            if not match:
                raise HTTPException(status_code=400, detail=f"Invalid location reference for text file: '{location}'")
            line_idx = int(match.group(1)) - 1
            
            with open(file_path, "r", encoding="utf-8") as f:
                lines = f.read().split("\n")
                
            if line_idx < 0 or line_idx >= len(lines):
                raise HTTPException(status_code=400, detail=f"Line index {line_idx+1} out of bounds.")
                
            lines[line_idx] = text
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file format: {ext}")
            
        # Save as new version
        import datetime
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        filename_orig = doc_meta.get("filename", "document")
        base, extension = os.path.splitext(filename_orig)
        
        base_clean = re.sub(r'_v_\d{8}_\d{6}$', '', base)
        new_filename = f"{base_clean}_v_{timestamp}{extension}"
        new_file_path = f"{UPLOAD_DIR}/{new_filename}"
        
        if ext == "docx":
            doc_obj.save(new_file_path)
        elif ext == "pptx":
            prs.save(new_file_path)
        elif ext == "xlsx":
            wb.save(new_file_path)
        elif ext == "pdf":
            doc_obj.save(new_file_path)
            doc_obj.close()
        elif ext == "txt":
            with open(new_file_path, "w", encoding="utf-8") as f:
                f.write("\n".join(lines))
                
        print(f"\n[SERVER CONFIRMATION] /edit/apply successfully wrote changes to {ext.upper()} file at '{new_file_path}' (location: '{location}')")
                
        # Parse new file content map
        new_extracted_text = ""
        new_content_map = {}
        new_file_id = uuid.uuid4().hex
        
        if ext == "pdf":
            new_doc = fitz.open(new_file_path)
            pages = []
            extracted_text_parts = []
            for page_idx, page in enumerate(new_doc):
                blocks = page.get_text("blocks")
                block_list = []
                for b in blocks:
                    block_list.append({
                        "text": b[4],
                        "bbox": [b[0], b[1], b[2], b[3]],
                        "block_number": b[5],
                        "block_type": b[6]
                    })
                pages.append({
                    "page_index": page_idx,
                    "width": page.rect.width,
                    "height": page.rect.height,
                    "blocks": block_list
                })
                extracted_text_parts.append(page.get_text())
            new_extracted_text = "\n".join(extracted_text_parts)
            new_content_map = {"pages": pages}
            new_doc.close()
            
        elif ext == "docx":
            new_doc = docx.Document(new_file_path)
            paragraphs = []
            extracted_text_parts = []
            for idx, para in enumerate(new_doc.paragraphs):
                runs_list = []
                for run in para.runs:
                    color_val = None
                    if run.font and run.font.color:
                        try:
                            if run.font.color.rgb:
                                color_val = str(run.font.color.rgb)
                        except Exception:
                            pass
                    runs_list.append({
                        "text": run.text,
                        "bold": run.bold,
                        "italic": run.italic,
                        "font_name": run.font.name if run.font else None,
                        "font_size": run.font.size.pt if (run.font and run.font.size) else None,
                        "color": color_val
                    })
                paragraphs.append({
                    "index": idx,
                    "text": para.text,
                    "runs": runs_list
                })
                extracted_text_parts.append(para.text)
            new_extracted_text = "\n".join(extracted_text_parts)
            new_content_map = {"paragraphs": paragraphs}
            
        elif ext == "xlsx":
            new_wb = openpyxl.load_workbook(new_file_path, data_only=True)
            sheets = []
            extracted_text_parts = []
            for sheet_name in new_wb.sheetnames:
                ws = new_wb[sheet_name]
                cells_list = []
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            val_str = str(cell.value)
                            extracted_text_parts.append(f"Cell {cell.coordinate}: {val_str}")
                            formatting = {}
                            if cell.font:
                                color_val = None
                                if cell.font.color and hasattr(cell.font.color, 'rgb'):
                                    color_val = str(cell.font.color.rgb)
                                formatting = {
                                    "bold": cell.font.bold,
                                    "italic": cell.font.italic,
                                    "font_name": cell.font.name,
                                    "font_size": cell.font.size,
                                    "color": color_val
                                }
                            if cell.number_format:
                                formatting["number_format"] = cell.number_format
                            cells_list.append({
                                "address": cell.coordinate,
                                "value": val_str,
                                "formatting": formatting
                            })
                sheets.append({
                    "sheet_name": sheet_name,
                    "cells": cells_list
                })
            new_extracted_text = "\n".join(extracted_text_parts)
            new_content_map = {"sheets": sheets}
            
        elif ext == "pptx":
            new_prs = Presentation(new_file_path)
            slides = []
            extracted_text_parts = []
            for slide_idx, slide in enumerate(new_prs.slides):
                shapes_list = []
                def crawl_shape(shape):
                    shape_data = {
                        "name": shape.name,
                        "type": str(shape.shape_type),
                        "text": None,
                        "formatting": {}
                    }
                    if shape.has_text_frame:
                        text_val = shape.text.strip()
                        if text_val:
                            shape_data["text"] = text_val
                            extracted_text_parts.append(text_val)
                        if shape.text_frame.paragraphs:
                            p = shape.text_frame.paragraphs[0]
                            if p.runs:
                                r = p.runs[0]
                                color_val = None
                                if r.font and r.font.color:
                                    try:
                                        if r.font.color.rgb:
                                            color_val = str(r.font.color.rgb)
                                    except Exception:
                                        pass
                                shape_data["formatting"] = {
                                    "bold": r.font.bold if r.font else None,
                                    "italic": r.font.italic if r.font else None,
                                    "font_name": r.font.name if r.font else None,
                                    "font_size": r.font.size.pt if (r.font and r.font.size) else None,
                                    "color": color_val
                                }
                    if shape.has_table:
                        table_cells = []
                        for row in shape.table.rows:
                            row_cells = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                row_cells.append(cell_text)
                                if cell_text:
                                    extracted_text_parts.append(cell_text)
                            table_cells.append(row_cells)
                        shape_data["text"] = json.dumps(table_cells)
                    shapes_list.append(shape_data)
                    if shape.shape_type == 6:
                        for subshape in shape.shapes:
                            crawl_shape(subshape)
                for shape in slide.shapes:
                    crawl_shape(shape)
                slides.append({
                    "slide_index": slide_idx,
                    "shapes": shapes_list
                })
            new_extracted_text = "\n".join(extracted_text_parts)
            new_content_map = {"slides": slides}
            
        elif ext == "txt":
            with open(new_file_path, "r", encoding="utf-8") as f:
                new_extracted_text = f.read()
            paragraphs_list = []
            for idx, line in enumerate(new_extracted_text.split("\n")):
                if line.strip():
                    paragraphs_list.append({
                        "index": idx,
                        "text": line,
                        "runs": []
                    })
            new_content_map = {"paragraphs": paragraphs_list}
            
        content_maps[new_file_id] = new_content_map
        detected_style = classify_document_style(new_extracted_text)
        
        # Preserve original insights template
        new_insights = doc_meta.get("insights", {
            "summary": "Updated document.",
            "risks": "N/A",
            "recommendations": "N/A",
            "proactive_suggestions": []
        })
        
        uploaded_documents.append({
            "file_id": new_file_id,
            "filename": new_filename,
            "content": new_extracted_text,
            "style": detected_style,
            "content_map": new_content_map,
            "file_path": new_file_path,
            "insights": new_insights
        })
        
        # Add to vector DB
        chunks = text_splitter.split_text(new_extracted_text)
        docs_to_insert = [Document(page_content=chunk, metadata={"filename": new_filename}) for chunk in chunks]
        if docs_to_insert:
            vectorstore.add_documents(docs_to_insert)
            
        download_link = f"http://localhost:8000/files/{new_filename}"
        style_label = DOCUMENT_STYLE_PROFILES.get(detected_style, DOCUMENT_STYLE_PROFILES["casual"])["label"]
        
        word_count = len(new_extracted_text.split())
        short_summary = {
            "file_type": ext,
            "word_count": word_count
        }
        if ext == "pdf":
            short_summary["page_count"] = len(new_content_map.get("pages", []))
        elif ext == "docx":
            short_summary["paragraph_count"] = len(new_content_map.get("paragraphs", []))
        elif ext == "xlsx":
            short_summary["sheet_count"] = len(new_content_map.get("sheets", []))
        elif ext == "pptx":
            short_summary["slide_count"] = len(new_content_map.get("slides", []))
            
        return {
            "message": "Successfully applied edit",
            "file_id": new_file_id,
            "download_link": download_link,
            "filename": new_filename,
            "short_summary": short_summary,
            "insights": new_insights,
            "url": download_link,
            "content": new_extracted_text,
            "content_map": new_content_map,
            "document_style": detected_style,
            "document_style_label": style_label
        }
    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error applying edit to file: {str(e)}")
